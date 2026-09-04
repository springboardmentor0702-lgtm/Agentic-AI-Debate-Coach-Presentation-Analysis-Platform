import re
from pptx import Presentation as PptxPresentation

# A presentation "performance" score is only meaningful when the learner supplied
# enough actual speech. Slide-content quality is assessed independently.
MIN_SPEECH_WORDS = 8


def parse_pptx(path):
    prs = PptxPresentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        title = texts[0] if texts else f"Slide {i}"
        notes = ""
        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass
        slides.append({
            "slide_number": i,
            "title": title[:200],
            "text": "\n".join(texts),
            "speaker_notes": notes,
            "shape_count": len(slide.shapes),
        })
    return slides


def _content_feedback(text):
    words = text.split()
    has_numbers = bool(re.search(r"\b\d+(?:\.\d+)?%?\b", text))
    has_evidence = bool(re.search(
        r"\b(study|research|data|survey|source|evidence|report|result|citation)\b",
        text, re.I
    ))
    has_structure = bool(re.search(
        r"\b(first|second|third|because|therefore|however|process|architecture|method|result|conclusion)\b",
        text, re.I
    ))
    score = 45
    if len(words) >= 8:
        score += 15
    if len(words) >= 25:
        score += 10
    if has_evidence:
        score += 12
    if has_numbers:
        score += 8
    if has_structure:
        score += 8
    score = min(score, 98)

    strengths = []
    if len(words) >= 8:
        strengths.append("The slide contains enough actual content to identify a main idea.")
    if has_evidence:
        strengths.append("Evidence-related content is explicitly present in the slide text.")
    if not strengths:
        strengths.append("The uploaded slide was parsed successfully.")

    weaknesses, suggestions = [], []
    if not has_evidence:
        weaknesses.append("No explicit evidence or source is visible in the slide text.")
        suggestions.append("Add one concrete source, result, example, or measurable fact where appropriate.")
    if len(words) < 8:
        weaknesses.append("The slide text is very brief, so the main claim may need spoken explanation.")
        suggestions.append("Use the slide as a visual anchor and state the key claim clearly while presenting.")
    if not suggestions:
        suggestions.append("State one primary message and connect it directly to the strongest evidence.")
    return score, strengths, weaknesses, suggestions


def _speech_metrics(speech):
    words = len(speech.split())
    if words < MIN_SPEECH_WORDS:
        return {
            "sufficient": False,
            "word_count": words,
            "duration": None,
            "pace": None,
            "fillers": None,
            "clarity": None,
            "confidence": None,
            "delivery": None,
            "status": "Insufficient speech for meaningful presentation analysis.",
        }

    # These are transcript-derived estimates, not microphone/audio measurements.
    duration = round(words / 2.3, 1)
    pace = round(words / (duration / 60), 1)
    fillers = len(re.findall(r"\b(um|uh|like|you know|actually|basically)\b", speech, re.I))
    hedges = len(re.findall(r"\b(maybe|perhaps|i think|i guess|sort of|kind of|might|could be)\b", speech, re.I))
    clarity = round(max(0, min(100, 92 - fillers * 2 - max(0, pace - 155) * 0.12)), 1)
    confidence = round(max(0, min(100, 82 - hedges * 2 - fillers * 1.5)), 1)
    delivery = round((clarity + confidence) / 2, 1)
    return {
        "sufficient": True,
        "word_count": words,
        "duration": duration,
        "pace": pace,
        "fillers": fillers,
        "clarity": clarity,
        "confidence": confidence,
        "delivery": delivery,
        "status": "Transcript-derived estimate from captured speech; not an audio/acoustic measurement.",
    }



def analyze_slides(slides, speeches=None, audio_metrics=None):
    speeches = speeches or {}
    audio_metrics = audio_metrics or {}
    out = []

    def has_real_audio(metrics):
        if not isinstance(metrics, dict):
            return False
        try:
            duration = float(metrics.get("duration_seconds") or 0)
            active = float(metrics.get("active_seconds") or 0)
            return duration > 0 and active > 0
        except (TypeError, ValueError):
            return False

    for s in slides:
        speech = str(
            speeches.get(str(s["slide_number"]), speeches.get(s["slide_number"], "")) or ""
        ).strip()

        am = audio_metrics.get(
            str(s["slide_number"]),
            audio_metrics.get(s["slide_number"], {})
        ) or {}

        audio_available = has_real_audio(am)
        sm = _speech_metrics(speech)

        content_score, strengths, weaknesses, suggestions = _content_feedback(s["text"])

        if audio_available and sm["sufficient"]:
            duration = float(am.get("duration_seconds") or sm["duration"])

            sm["duration"] = round(duration, 2)
            sm["pace"] = (
                round(sm["word_count"] / (duration / 60), 1)
                if duration > 0 else None
            )

            actual_speech = speech
            explanation = round(
                content_score * 0.65 + sm["delivery"] * 0.35, 1
            )
            text_clarity = sm["clarity"]
            spoken_clarity = sm["clarity"]
            confidence = sm["confidence"]
            delivery_quality = sm["delivery"]
            speech_status = "Recorded microphone speech analyzed with transcript evidence."

        elif speech:
            actual_speech = speech
            explanation = None
            text_clarity = _text_clarity(speech)
            spoken_clarity = None
            confidence = None
            delivery_quality = None
            speech_status = (
                "Typed/transcribed text is available, but no microphone recording "
                "was provided. Spoken performance metrics are unavailable."
            )

            strengths = list(strengths) + [
                "The provided text can be reviewed for wording and clarity."
            ]
            weaknesses = list(weaknesses) + [
                "No microphone recording was provided, so spoken confidence, "
                "communication and delivery cannot be measured."
            ]
            suggestions = list(suggestions) + [
                "Record your actual presentation to measure confidence, "
                "communication, pacing and delivery."
            ]

        else:
            actual_speech = None
            explanation = None
            text_clarity = None
            spoken_clarity = None
            confidence = None
            delivery_quality = None
            speech_status = "Unavailable ? no speech was captured for this slide."

        out.append({
            **s,
            "speech": actual_speech,
            "speech_captured": bool(speech),
            "speech_sufficient": bool(audio_available and sm["sufficient"]),
            "audio_available": audio_available,
            "speech_word_count": sm["word_count"],
            "speech_duration_seconds": sm["duration"] if audio_available else None,
            "audio_metrics": am if am else None,
            "speech_pace_wpm": sm["pace"] if audio_available else None,
            "filler_words": sm["fillers"] if audio_available else None,
            "clarity": spoken_clarity,
            "text_clarity": text_clarity,
            "confidence_estimate": confidence,
            "explanation_quality": explanation,
            "content_quality": content_score,
            "relevance": None,
            "delivery_quality": delivery_quality,
            "engagement_indicators": speech_status,
            "what_you_said": actual_speech,
            "strengths": strengths,
            "weaknesses": weaknesses,
            "suggestions": suggestions,
        })

    content_scores = [
        x["content_quality"] for x in out
        if x["content_quality"] is not None
    ]

    performance_slides = [
        x for x in out
        if x["audio_available"] and x["speech_sufficient"]
    ]

    delivery = [
        x["delivery_quality"] for x in performance_slides
        if x["delivery_quality"] is not None
    ]

    explanations = [
        x["explanation_quality"] for x in performance_slides
        if x["explanation_quality"] is not None
    ]

    text_clarity_values = [
        x["text_clarity"] for x in out
        if x["text_clarity"] is not None
    ]

    confidence_values = [
        x["confidence_estimate"] for x in performance_slides
        if x["confidence_estimate"] is not None
    ]

    any_audio = any(x["audio_available"] for x in out)

    overall = (
        round(sum(explanations) / len(explanations), 1)
        if explanations
        else (
            round(sum(content_scores) / len(content_scores), 1)
            if content_scores else None
        )
    )

    return {
        "slides": out,
        "performance_available": bool(performance_slides),
        "audio_available": any_audio,
        "overall_score": overall,
        "content_quality": (
            round(sum(content_scores) / len(content_scores), 1)
            if content_scores else None
        ),
        "explanation_quality": (
            round(sum(explanations) / len(explanations), 1)
            if explanations else None
        ),
        "speech_pace": (
            "Observed microphone duration + transcript word count"
            if delivery
            else "Unavailable ? record your presentation to measure spoken pace."
        ),
        "filler_words": (
            sum((x["filler_words"] or 0) for x in performance_slides)
            if performance_slides else None
        ),
        "confidence": (
            round(sum(confidence_values) / len(confidence_values), 1)
            if confidence_values else None
        ),
        "clarity": (
            round(
                sum(x["clarity"] for x in performance_slides if x["clarity"] is not None)
                / len([x for x in performance_slides if x["clarity"] is not None]),
                1
            )
            if any(x["clarity"] is not None for x in performance_slides)
            else None
        ),
        "text_clarity": (
            round(sum(text_clarity_values) / len(text_clarity_values), 1)
            if text_clarity_values else None
        ),
        "engagement": (
            "Browser microphone activity estimate"
            if delivery
            else "Unavailable ? microphone recording required."
        ),
        "transitions": None,
        "structure": (
            round(sum(content_scores) / len(content_scores), 1)
            if content_scores else None
        ),
        "delivery": (
            round(sum(delivery) / len(delivery), 1)
            if delivery else None
        ),
        "best_slide": (
            max(out, key=lambda x: x["content_quality"])["slide_number"]
            if out else None
        ),
        "weakest_slide": (
            min(out, key=lambda x: x["content_quality"])["slide_number"]
            if out else None
        ),
        "most_improved_slide": None,
        "biggest_delivery_issue": (
            "No microphone recording was provided, so delivery cannot be measured."
            if not delivery
            else
            "Review pace, pauses and filler words on slides with the lowest delivery estimate."
        ),
        "top_recommendations": [
            (
                "Record your actual presentation to measure confidence, "
                "communication, pacing and delivery."
                if not any_audio
                else
                "Keep using recorded speech so delivery feedback is grounded in microphone evidence."
            ),
            "Keep one primary message per slide.",
            "Tie slide claims to concrete evidence.",
        ],
    }


def _text_clarity(text):
    words = text.split()
    if not words:
        return None

    sentences = max(1, len(re.findall(r"[.!?]+", text)))
    avg_sentence = len(words) / sentences

    score = 92
    if avg_sentence > 28:
        score -= 10
    elif avg_sentence > 20:
        score -= 5

    fillers = len(
        re.findall(r"\b(um|uh|like|you know|actually|basically)\b", text, re.I)
    )
    score -= min(15, fillers * 2)

    return round(max(0, min(100, score)), 1)

