import os,sys,tempfile
sys.path.insert(0,os.path.abspath(os.path.join(os.path.dirname(__file__),"..")))
from pptx import Presentation
from app.services.presentation import parse_pptx,analyze_slides

def test_pptx_parse_and_analysis():
    prs=Presentation(); slide=prs.slides.add_slide(prs.slide_layouts[1]);slide.shapes.title.text='Evidence';slide.placeholders[1].text='Use credible studies and explain the warrant.'
    with tempfile.NamedTemporaryFile(suffix='.pptx') as f:
        prs.save(f.name); slides=parse_pptx(f.name)
    assert len(slides)==1 and slides[0]['title']=='Evidence'
    out=analyze_slides(slides,{'1':'Research shows strong outcomes. Um, the evidence supports the claim.'})
    assert out['overall_score']>0 and out['slides'][0]['filler_words']>=1
